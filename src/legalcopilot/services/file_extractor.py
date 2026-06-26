"""File text extraction — extract readable text from uploaded documents.

Supports: PDF, DOCX, XLSX, PPTX, TXT.
Falls back gracefully when extraction libraries are missing.
"""

import io
import logging
import os

logger = logging.getLogger(__name__)


def extract_text(file_bytes: bytes, filename: str) -> str:
    """Extract text content from file bytes based on file extension.

    Args:
        file_bytes: Raw file content.
        filename: Original filename (used to determine type by extension).

    Returns:
        Extracted text string. Empty string if extraction fails or is unsupported.
    """
    ext = os.path.splitext(filename)[1].lower()

    extractors = {
        ".pdf": _extract_pdf,
        ".docx": _extract_docx,
        ".xlsx": _extract_xlsx,
        ".pptx": _extract_pptx,
        ".txt": _extract_txt,
    }

    extractor = extractors.get(ext)
    if extractor is None:
        logger.warning("No extractor for extension: %s", ext)
        return ""

    try:
        return extractor(file_bytes)
    except Exception:
        logger.exception("Text extraction failed for %s", filename)
        return ""


_MAX_PDF_PAGES = 500
_MAX_XLSX_ROWS = 10_000
_MAX_PPTX_SLIDES = 200
_MAX_TEXT_BYTES = 10_000_000  # 10 MB


def _extract_pdf(file_bytes: bytes) -> str:
    """Extract text from PDF using pypdf."""
    from pypdf import PdfReader

    reader = PdfReader(io.BytesIO(file_bytes))
    pages = []
    for i, page in enumerate(reader.pages):
        if i >= _MAX_PDF_PAGES:
            pages.append(f"[Truncated: {len(reader.pages) - _MAX_PDF_PAGES} pages omitted]")
            break
        text = page.extract_text()
        if text:
            pages.append(text)
    return "\n\n".join(pages)


def _extract_docx(file_bytes: bytes) -> str:
    """Extract text from DOCX using python-docx."""
    from docx import Document

    doc = Document(io.BytesIO(file_bytes))
    paragraphs = []
    for para in doc.paragraphs:
        if para.text.strip():
            paragraphs.append(para.text)
    return "\n\n".join(paragraphs)


def _extract_xlsx(file_bytes: bytes) -> str:
    """Extract text from XLSX using openpyxl."""
    from openpyxl import load_workbook

    wb = load_workbook(io.BytesIO(file_bytes), read_only=True, data_only=True)
    sheets = []
    total_rows = 0
    for ws in wb.worksheets:
        rows = []
        for row in ws.iter_rows(values_only=True):
            if total_rows >= _MAX_XLSX_ROWS:
                rows.append(f"[Truncated: row limit ({_MAX_XLSX_ROWS}) reached]")
                break
            cells = [str(c) if c is not None else "" for c in row]
            if any(cells):
                rows.append(" | ".join(cells))
                total_rows += 1
        if rows:
            sheets.append(f"[Sheet: {ws.title}]\n" + "\n".join(rows))
        if total_rows >= _MAX_XLSX_ROWS:
            break
    wb.close()
    return "\n\n".join(sheets)


def _extract_pptx(file_bytes: bytes) -> str:
    """Extract text from PPTX using python-pptx."""
    from pptx import Presentation

    prs = Presentation(io.BytesIO(file_bytes))
    slides = []
    for i, slide in enumerate(prs.slides[:_MAX_PPTX_SLIDES], 1):
        texts = []
        for shape in slide.shapes:
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    if para.text.strip():
                        texts.append(para.text)
        if texts:
            slides.append(f"[Slide {i}]\n" + "\n".join(texts))
    return "\n\n".join(slides)


def _extract_txt(file_bytes: bytes) -> str:
    """Extract text from plain text file."""
    if len(file_bytes) > _MAX_TEXT_BYTES:
        file_bytes = file_bytes[:_MAX_TEXT_BYTES]
    try:
        return file_bytes.decode("utf-8")
    except UnicodeDecodeError:
        return file_bytes.decode("latin-1")
